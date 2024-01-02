from typing import List, Optional, Iterator

from langchain.docstore.document import Document
from langchain.document_loaders.base import BaseLoader
from pptx import Presentation


class PPTLoader(BaseLoader):
    """
    Load a document from an Excel file.
    """
    def __init__(self, path: str, sheet_name: Optional[str] = None, *args, **kwargs):
        self.path = path

    def load(self) -> List[Document]:
        return list(self.lazy_load())

    def lazy_load(self) -> Iterator[Document]:
        pr = Presentation(self.file_path)

        for slide in pr.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    yield   Document(
                            page_content=shape.text, metadata={"source": slide.slide_id}
                        )
